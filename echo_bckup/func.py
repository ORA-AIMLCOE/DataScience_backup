import io
import pdfplumber
import pandas as pd
import re
from pptx import Presentation

def extract_text_from_pptx(pptx_path):
    prs = Presentation(pptx_path)
    text_content = []
    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                text_content.append(shape.text)
    return "\n".join(text_content)

def extract_section(text, section_name, stop_keywords):
    section_content = []
    lines = text.split("\n")
    capture = False

    for line in lines:
        if section_name.lower() in line.lower():
            capture = True
            section_content.append(line)
        elif capture:
            if any(stop_keyword.lower() in line.lower() for stop_keyword in stop_keywords) or \
               line.strip().isdigit():
                break
            section_content.append(line)

    cleaned_section = re.sub(r'[\x0b\x0c]', '', "\n".join(section_content)).strip()
    cleaned_section = cleaned_section.replace("ScopeEnvironment setup", "").strip()
    return cleaned_section

def handler(ctx, data: io.BytesIO = None):
    try:
        # Load PDF from input (assuming it comes as binary)
        with pdfplumber.open("/function/rfp_sample.pdf") as pdf:
            sections = {}
            current_section = None
            section_text = ""

            for page in pdf.pages:
                text = page.extract_text()
                lines = text.split("\n")

                for line in lines:
                    if line.isupper() or (line[0].isdigit() and line[1:2] == "."):
                        if current_section:
                            sections[current_section] = section_text.strip()
                        current_section = line
                        section_text = ""
                    else:
                        section_text += line + " "

            if current_section:
                sections[current_section] = section_text.strip()

        sections_df = pd.DataFrame(list(sections.items()), columns=["Section", "Content"])

        # Extract text from PPTX (assuming file exists in function container)
        pptx_text = extract_text_from_pptx("/function/rfp_response.pptx")
        stop_keywords = ["Copyright", "Scope", "Business Value", "Data Sources", "Deliverables", "Assumptions"]
        environment_setup_section = extract_section(pptx_text, "Environment setup", stop_keywords)

        # Extract offer names from environment setup section
        offer_names = [line.strip() for line in environment_setup_section.split("\n") if line.strip() and "setup" not in line.lower()]

        # Create BOM DataFrame
        data = []
        for offer_name in offer_names:
            data.append({"License Included PaaS": "Prod", "Offer Name": offer_name})
            data.append({"License Included PaaS": "Non-Prod", "Offer Name": offer_name})

        df_bom = pd.DataFrame(data)

        # Save the BOM to a CSV file inside the function's filesystem
        output_path = "/tmp/function/bom_with_offer_names.csv"
        df_bom.to_csv(output_path, index=False)

        return {"status": "success", "message": f"BOM generated and saved to {output_path}"}

    except Exception as e:
        return {"status": "error", "message": str(e)}
